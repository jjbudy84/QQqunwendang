from __future__ import annotations

from pathlib import Path


def parse_document(path: Path, original_name: str | None = None) -> str:
    suffix = path.suffix.lower() or Path(original_name or "").suffix.lower()
    if suffix == ".txt":
        return parse_txt(path)
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix in {".ppt", ".pptx"}:
        return parse_pptx(path)
    if suffix in {".doc", ".docx"}:
        return parse_docx(path)
    raise ValueError(f"不支持的文件类型：{suffix}")


def parse_txt(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def parse_pdf(path: Path) -> str:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[第 {index} 页]\n{text}")
    return "\n\n".join(pages)


def parse_pptx(path: Path) -> str:
    if path.suffix.lower() == ".ppt":
        raise ValueError("旧版 .ppt 需要先另存为 .pptx，或安装 LibreOffice 后扩展转换逻辑。")

    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
        if lines:
            slides.append(f"[第 {slide_index} 页]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def parse_docx(path: Path) -> str:
    if path.suffix.lower() == ".doc":
        raise ValueError("旧版 .doc 需要先另存为 .docx，或安装 antiword/LibreOffice 后扩展转换逻辑。")

    import docx

    document = docx.Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]

    table_lines: list[str] = []
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_lines.append(" | ".join(cells))

    return "\n".join(paragraphs + table_lines)
